# -*- coding: utf-8 -*-
import os
import glob
import argparse
import logging
from datetime import datetime
from logging import getLogger, StreamHandler, FileHandler, Formatter
from pptx import Presentation
from pptx.util import Cm
import pandas as pd
import numpy as np

def initialize_logger(log_file):
    logger = getLogger()
    logger.setLevel(logging.INFO)

    file_handler = FileHandler(log_file)
    stream_handler = StreamHandler()

    file_handler.setLevel(logging.INFO)
    stream_handler.setLevel(logging.INFO)
    handler_formatter = Formatter('%(message)s')
    file_handler.setFormatter(handler_formatter)
    stream_handler.setFormatter(handler_formatter)

    logger.addHandler(file_handler)
    logger.addHandler(stream_handler)

    return logger


def print_log(msg):
    if 'LOGGER' in globals():
        LOGGER.info(msg)
    else:
        print(msg)


def get_tables_in_slide(slide):
    # https://python-pptx.readthedocs.io/en/latest/dev/analysis/shp-graphfrm.html
    shapes = slide.shapes
    table_list = [shp.table for shp in shapes if shp.has_table]
    return table_list


def get_table_label(table):
    # Returns the content of the top-left cell that I decided to label the table
    label = table.rows[1].cells[0].text
    return label


def map_table_label(table_list):
    # Returns a dict of {label: table}
    # notice:The label is capitalized
    table_map = {}
    for tbl in table_list:
        label = get_table_label(tbl)
        table_map[label.strip().upper()] = tbl
    return table_map


def set_dataframe_to_table(df, tbl):
    # set the items of the dataframe to the cells of the table
    rows_iter = iter(tbl.rows)
    next(rows_iter) # skip the header row
    for row_idx, row in enumerate(rows_iter):
        if df.shape[0] <= row_idx:
            break
        for col_idx, cell in enumerate(row.cells):
            item = df.iloc[row_idx, col_idx]
            try:
                str_value = '{:.3g}'.format(item)
            except ValueError:
                str_value = str(item)
            # Notice: All cells in the template table cannot be empty to have at least one 'run'.
            runs = cell.text_frame.paragraphs[0].runs
            runs_num = len(runs)
            if runs_num > 1:
                runs[0].text = str_value
                # clear the redundant runs
                for i in range(1, runs_num):
                    runs[i].text = ''
            else:
                runs[0].text = str_value


def select_wanted_path(summary_dir, part_of_filename):
    # Find the only one wanted filepath from similar ones
    candidate_str = os.path.join(summary_dir, f'{part_of_filename}*.jpg')
    candidates = glob.glob(candidate_str)
    selected_paths = [c for c in candidates if 'logodds' not in c]
    try:
        slected_path = selected_paths[0]
    except IndexError:
        print(f'No match {candidate_str}')
    return slected_path


def select_wanted_path_list(summary_dir, part_of_filename):
    """ Find the mutliple wanted filepaths from similar ones
    """
    candidate_str = os.path.join(summary_dir, f'{part_of_filename}*.jpg')
    candidates = sorted(glob.glob(candidate_str))
    selected_paths = [c for c in candidates if 'logodds' not in c]
    try:
        _ = selected_paths[0]
    except IndexError:
        print_log(f'No match {candidate_str}')
    return selected_paths


def get_text_frames_in_slide(slide, tag_str=''):
    # returns the text frames in the slide
    # returns the text frames with the specified str
    shapes = slide.shapes
    text_frames = [shp.text_frame for shp in shapes if shp.has_text_frame]

    if tag_str != '':
        text_frames = [
            tf for tf in text_frames if tf.paragraphs[0].runs[0].text.strip().upper() == tag_str]

    return text_frames


def make_summary_label(summary_dir):
    """Returns a string concatinating experiment labels summaried in the summary dir

    Args:
        summary_dir (str): A path to the summary directory

    Returns:
        str: a string for the summary label
    """
    stage_time_df = pd.read_csv(os.path.join(
        summary_dir, 'stage-time_table.csv'))
    sum_label = '_'.join(stage_time_df.iloc[:, 0].unique())
    return sum_label


def assign_page_title(prs, title):
    """ Assign the page title to all slides

    Args:
        prs (Presentation): The slide object
        title (str): The slide title to be inserted
    """
    for si in range(7):
        slide = prs.slides[si]
        txt_frms = get_text_frames_in_slide(slide, 'SUMMARY LABEL')
        try:
            txt_frms[0].paragraphs[0].runs[0].text = title
        except IndexError:
            print(f'Failed to set the experiment label to the slide #{si}')


def prep_table_of_stage_stats(prs, summary_dir):
    """Prepare tables of the stage stats

    Args:
        prs (Presentation): The slide object
        summary_dir (str): The path to the summary directory
    """
    df_stage_time = pd.read_csv(os.path.join(
        summary_dir, 'stage-time_stats_table.csv'))
    df_sw_trans = pd.read_csv(os.path.join(
        summary_dir, 'sleep-wake-transition_probability_stats_table.csv'))
    df_stage_trans = pd.read_csv(os.path.join(
        summary_dir, 'stage-transition_probability_stats_table.csv'))

    slide = prs.slides[0]
    table_list = get_tables_in_slide(slide)
    table_map = map_table_label(table_list)

    set_dataframe_to_table(df_stage_time.fillna(' '), table_map['STAGE-TIME'])
    set_dataframe_to_table(df_sw_trans.fillna(' '), table_map['SW-TRANS'])
    set_dataframe_to_table(df_stage_trans.fillna(' '), table_map['STAGE-TRANS'])


def prep_table_of_dpd(prs, summary_dir):
    """Prepare tables of the delta-power dynamics

    Args:
        prs (Presentation): The slide object
        summary_dir (str): The path to the summary directory
    """

    dpd_path = os.path.join(summary_dir, 'delta_power_dynamics')
    if os.path.exists(dpd_path) is False:
        print_log('No delta power analysis data found')
        return 

    df_dpd_stats = pd.read_csv(os.path.join(
        dpd_path, 'delta-power-dynamics_stats_table.csv'))

    slide = prs.slides[8]
    table_list = get_tables_in_slide(slide)
    table_map = map_table_label(table_list)

    set_dataframe_to_table(df_dpd_stats.fillna(' '), table_map['DPD-STATS'])


def prep_table_of_psd(prs, summary_dir):
    """Prepare tables of the PDS stats

    Args:
        prs (Presentation): The slide object
        summary_dir (str): The path to the summary directory
    """
    time_domains = np.tile(['allday', 'first-halfday', 'second-halfday'], 2)
    transform_types = np.repeat(['linear', 'log'], 3)
    slide_idx = [2, 3, 4, 5, 6, 7]

    for si, tf, td in zip(slide_idx, transform_types, time_domains):
        df_psd_ff = pd.read_csv(os.path.join(
            summary_dir, 'PSD_raw', f'PSD_raw_none_{tf}_{td}_profile_stats_table.csv'))
        df_psd_ft = pd.read_csv(os.path.join(
            summary_dir, 'PSD_norm', f'PSD_norm_none_{tf}_{td}_profile_stats_table.csv'))
        df_psd_tf = pd.read_csv(os.path.join(
            summary_dir, 'PSD_raw', f'PSD_raw_TDD_{tf}_{td}_profile_stats_table.csv'))
        df_psd_tt = pd.read_csv(os.path.join(
            summary_dir, 'PSD_norm', f'PSD_norm_AUC_{tf}_{td}_profile_stats_table.csv'))

        slide = prs.slides[si]
        table_list = get_tables_in_slide(slide)
        table_map = map_table_label(table_list)

        set_dataframe_to_table(
            df_psd_ff.iloc[16:20, :].fillna(' '), table_map['PSD-FF'])
        set_dataframe_to_table(
            df_psd_ft.iloc[16:20, :].fillna(' '), table_map['PSD-FT'])
        set_dataframe_to_table(
            df_psd_tf.iloc[16:20, :].fillna(' '), table_map['PSD-TF'])
        set_dataframe_to_table(
            df_psd_tt.iloc[16:20, :].fillna(' '), table_map['PSD-TT'])


def prep_fig_of_stage_stats(prs, summary_dir):
    """Prepare plots of stage stats (page 1)

    Args:
        prs (Presentation): The slide object
        summary_dir (str): The path to the summary directory
    """
    slide = prs.slides[0]

    path_stage_time = os.path.join(summary_dir, f'stage-time_barchart.jpg')
    path_sw_trans_bar = select_wanted_path(
        summary_dir, 'sleep-wake-transition_probability_barchart')
    path_sw_trans_circ = select_wanted_path(
        summary_dir, 'sleep-wake-transition_circadian_profile_G')
    path_stage_prof = select_wanted_path(summary_dir, 'stage-time_profile_G')
    path_stage_trans_bar = select_wanted_path(
        summary_dir, 'stage-transition_probability_barchart')

    slide.shapes.add_picture(path_stage_time, Cm(0.95),
                             Cm(3.4), Cm(11.05), Cm(5))
    slide.shapes.add_picture(
        path_sw_trans_bar, Cm(14.2), Cm(3.4), Cm(5.3), Cm(5))
    slide.shapes.add_picture(path_sw_trans_circ, Cm(
        19.6), Cm(3.4), Cm(13.6), Cm(5))
    slide.shapes.add_picture(path_stage_prof, Cm(
        0.42), Cm(11.42), Cm(14.52), Cm(7.6))
    slide.shapes.add_picture(path_stage_trans_bar, Cm(
        14.74), Cm(10.73), Cm(11.34), Cm(8.1))


def prep_fig_of_dpd(prs, summary_dir):
    """Prepare plots of the delta-power dynamics (page 6)

    Args:
        prs (Presentation): The slide object
        summary_dir (str): The path to the summary directory
    """
    slide = prs.slides[5]
    dpd_path = os.path.join(summary_dir, 'delta_power_dynamics')

    if os.path.exists(dpd_path) is False:
        # 'No delta power analysis data found' is printed in the prep_table_of_dpd function
        return 

    # delta-power dyanmics simulation (group each)
    path_dpd_sim_ge_list = select_wanted_path_list(dpd_path, 'delta-power-dynamics_GE')

    slide.shapes.add_picture(path_dpd_sim_ge_list[0], Cm(0.7),
                             Cm(3.0), Cm(16), Cm(4.74))

    # Asymptotes boxchart
    path_asymp_boxplot = select_wanted_path(dpd_path, 'delta-power-dynamics_comparison_of_asymptotes')

    try:
        slide.shapes.add_picture(path_dpd_sim_ge_list[1], Cm(0.7),
                                Cm(8.1), Cm(16), Cm(4.74))

        # delta-power dynamics simulation (group comparison)
        path_dpd_sim = select_wanted_path(dpd_path, 'delta-power-dynamics_GC')
        slide.shapes.add_picture(path_dpd_sim, Cm(0.7),
                                Cm(13.8), Cm(16), Cm(4.74))

        # Taus 2D plot
        path_taus_2d_plot = select_wanted_path(dpd_path, 'delta-power-dynamics_taus_2D-plot_GC')

        # Taus barchart
        path_taus_barchart = select_wanted_path(dpd_path, 'delta-power-dynamics_taus_barchart_GC')


    except IndexError:
        print_log('Only a single group found for the delta power dynamics simulation')

        # Taus 2D plot
        path_taus_2d_plot = select_wanted_path(dpd_path, 'delta-power-dynamics_taus_2D-plot_all-group')

        # Taus barchart
        path_taus_barchart = select_wanted_path(dpd_path, 'delta-power-dynamics_taus_barchart_all-group')


    slide.shapes.add_picture(path_taus_2d_plot, Cm(17.38),
                             Cm(0.62), Cm(8.0), Cm(7.48))

    slide.shapes.add_picture(path_taus_barchart, Cm(25.62),
                             Cm(3.01), Cm(8.0), Cm(5.13))

    slide.shapes.add_picture(path_asymp_boxplot, Cm(17.00),
                             Cm(13.77), Cm(7.41), Cm(5.13))


def prep_fig_of_power_timeseries(prs, summary_dir):
    """Prepare plots of power timeseries (page 2)

    Args:
        prs (Presentation): The slide object
        summary_dir (str): The path to the summary directory
    """
    slide = prs.slides[1]

    # (voltage normalization, spectrum normalization, left, top)
    parm_list = [('raw', 'none', 2.01, 3.43),
                 ('raw', 'TDD', 2.01, 11.5),
                 ('norm', 'none', 17.64, 3.43),
                 ('norm', 'AUC', 17.64, 11.5)]

    for vol_dist_type, scaling_type, left, top in parm_list:
        path_pwr_ts_delta_wake = select_wanted_path(os.path.join(
            summary_dir, f'PSD_{vol_dist_type}'), f'power-timeseries_{vol_dist_type}_{scaling_type}_linear_delta_Wake_G')
        path_pwr_ts_delta_nrem = select_wanted_path(os.path.join(
            summary_dir, f'PSD_{vol_dist_type}'), f'power-timeseries_{vol_dist_type}_{scaling_type}_linear_delta_NREM_G')
        path_pwr_ts_delta_all = select_wanted_path(os.path.join(
            summary_dir, f'PSD_{vol_dist_type}'), f'power-timeseries_{vol_dist_type}_{scaling_type}_linear_delta_G')

        picture1 = slide.shapes.add_picture(
            path_pwr_ts_delta_wake, Cm(left), Cm(top), Cm(15.26), Cm(2.44))
        picture2 = slide.shapes.add_picture(
            path_pwr_ts_delta_nrem, Cm(left), Cm(top+2.47), Cm(15.26), Cm(2.44))
        picture3 = slide.shapes.add_picture(path_pwr_ts_delta_all, Cm(
            left), Cm(top+2.47*2), Cm(15.26), Cm(2.44))

        # put the pictures at the bottom of z-order
        # https://stackoverflow.com/questions/58601247/how-to-set-picture-to-bottom-layer-in-ppt
        # pylint: disable = protected-access
        slide.shapes._spTree.insert(2, picture1._element)
        slide.shapes._spTree.insert(2, picture2._element)
        slide.shapes._spTree.insert(2, picture3._element)


def prep_fig_of_psd(prs, summary_dir):
    """Prepare plots of PSD (page 3,4,5)
    Args:
        prs (Presentation): The slide object
        summary_dir (str): The path to the summary directory
    """
    time_domains = np.tile(['allday', 'first-halfday', 'second-halfday'], 2)
    transform_types = np.repeat(['linear', 'log'], 3)

    slide_idx = [2, 3, 4, 5, 6, 7]
    # voltage normalization, spectrum normalization, left, top
    parm_list = [('raw', 'none', 1.98, 3.94),
                 ('norm', 'none', 18.2, 3.94),
                 ('raw', 'TDD', 1.98, 12.23),
                 ('norm', 'AUC', 18.2, 12.23)]

    for si, tf, td in zip(slide_idx, transform_types, time_domains):
        slide = prs.slides[si]
        for vol_dist_type, scaling_type, left, top in parm_list:
            path_psd = select_wanted_path(os.path.join(
                summary_dir, f'PSD_{vol_dist_type}'), f'PSD_{vol_dist_type}_{scaling_type}_{tf}_{td}_profile_G')
            slide.shapes.add_picture(path_psd, Cm(
                left), Cm(top), Cm(14.84), Cm(4.5))


def prep_fig_of_spindle(prs, summary_dir):
    """Prepare plots of spindle (page 7)

    Args:
        prs (Presentation): The slide object
        summary_dir (str): The path to the summary directory
    """
    slide = prs.slides[6]

    if os.path.exists(os.path.join(summary_dir, 'spindle')) is False:
        print_log('No spindle data found')
        return 

    path_spindle_density_in_NREM              = os.path.join(summary_dir, 'spindle', f'spindle_density_in_NREM.jpg')
    path_spindle_distributions_NREMbeforeREM  = os.path.join(summary_dir, 'spindle', 'spindle_distributions_NREMbeforeREM.jpg')
    path_spindle_distributions_NREMafterWAKE  = os.path.join(summary_dir, 'spindle', 'spindle_distributions_NREMafterWAKE.jpg')
    path_spindle_distributions_NREMbeforeWAKE = os.path.join(summary_dir, 'spindle', 'spindle_distributions_NREMbeforeWAKE.jpg')
    path_spindle_distributions_NREMafterREM   = os.path.join(summary_dir, 'spindle', 'spindle_distributions_NREMafterREM.jpg')

    slide.shapes.add_picture(path_spindle_density_in_NREM, Cm(2.28), Cm(3.43), Cm(6.97), Cm(10))
    slide.shapes.add_picture(path_spindle_distributions_NREMbeforeREM, Cm(12.65), Cm(3.67), Cm(8.28), Cm(7))
    slide.shapes.add_picture(path_spindle_distributions_NREMafterREM, Cm(24.35), Cm(3.55), Cm(8.28), Cm(7))
    slide.shapes.add_picture(path_spindle_distributions_NREMbeforeWAKE, Cm(12.65), Cm(12.05), Cm(8.28), Cm(7))
    slide.shapes.add_picture(path_spindle_distributions_NREMafterWAKE, Cm(24.35), Cm(12.03), Cm(8.28), Cm(7))


def make_slide(args):
    """The main function to make the slides

    Args:
        args : The command line arguments
    """
    summary_dir = args.summary_dir

    print_log(f'Making summary slides of {summary_dir}')
    path2template = os.path.join(os.path.dirname(__file__), r'faster2lib/EEG_power_specrum_template.pptx')
    print_log(f'The template pptx:{path2template}')

    prs = Presentation(path2template)

    # Assign the page title to all slides
    sum_label = make_summary_label(summary_dir)
    assign_page_title(prs, sum_label)

    # Prepare the stage stats (page 1)
    prep_table_of_stage_stats(prs, summary_dir)

    # Prepare tables of the PDS stats (page 3,4,5,6,7,8)
    prep_table_of_psd(prs, summary_dir)

    # Prepare the table of delta-power dynamics (page 9)
    prep_table_of_dpd(prs, summary_dir)

    # Prepare plots of stage stats (page 1)
    prep_fig_of_stage_stats(prs, summary_dir)

    # Prepare plots of power timeseries (page 2)
    prep_fig_of_power_timeseries(prs, summary_dir)

    # Prepare plots of PSD (page 3,4,5,6,7,8)
    prep_fig_of_psd(prs, summary_dir)

    # Prepare plots of delta-power dynamics (page 9)
    prep_fig_of_dpd(prs, summary_dir)

    # Prepare plots of spindel (page 10)
    prep_fig_of_spindle(prs, summary_dir)

    path2summary_slide = os.path.join(summary_dir, 'summary.pptx')
    prs.save(path2summary_slide)
    print_log(f'The summary slide is {path2summary_slide}')


if __name__ == '__main__':

    PARSER = argparse.ArgumentParser()
    PARSER.add_argument("-s", "--summary_dir", required=True,
                        help="paths to the summary directory")

    ARGS = PARSER.parse_args()

    DT_STR = datetime.now().strftime('%Y-%m-%d_%H%M%S')
    LOGGER = initialize_logger(os.path.join(
        ARGS.summary_dir, 'log', f'summary.{DT_STR}.log'))

    make_slide(ARGS)
